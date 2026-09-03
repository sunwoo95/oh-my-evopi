#!/usr/bin/env python3
"""evopi 팀 세미나 덱 생성기.

참조 시나리오: /opt/workspace/local/sw4kim/my-agent/D1-01.ClaudeCode-Architecture.pdf (45p, 10섹션)
테마: 그래파이트 다크 + 틸 액센트 (AWS 테마 비사용).
실행: /tmp/pptx-venv/bin/python build_deck.py  → evopi-architecture.pptx
"""
from __future__ import annotations

import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
DIAG = os.path.normpath(os.path.join(HERE, "..", "diagrams"))
OUT = os.path.join(HERE, "evopi-architecture.pptx")

# ── 테마 ────────────────────────────────────────────────────────────────────
BG = RGBColor(0x14, 0x17, 0x1C)
PANEL = RGBColor(0x1C, 0x20, 0x28)
CARD = RGBColor(0x22, 0x27, 0x33)
BORDER = RGBColor(0x34, 0x3A, 0x46)
ACCENT = RGBColor(0x2E, 0xC4, 0xB6)
ACCENT_DIM = RGBColor(0x1F, 0x6F, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xAA, 0xB0, 0xBB)
DIM = RGBColor(0x6F, 0x75, 0x80)
WARN = RGBColor(0xF0, 0x8A, 0x4B)
FONT = "Noto Sans CJK KR"
MONO = "DejaVu Sans Mono"

W, H = Inches(13.333), Inches(7.5)
LM = Inches(0.6)          # 좌 여백
CW = W - 2 * LM           # 콘텐츠 폭
FOOTER = "oh-my-evopi · 팀 세미나 · 2026-09"

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]
page_no = 0


# ── 저수준 헬퍼 ─────────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill=CARD, line=BORDER, radius=True, lw=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if radius:
        shp.adjustments[0] = 0.06
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    shp.text_frame.text = ""
    return shp


def text(slide, x, y, w, h, s, size=12, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, font=FONT, spacing=None):
    """s: str 또는 [(str, {size,color,bold,font}) ...] 런 리스트 / 줄은 '\n' 로 구분(각 줄이 문단)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    runs = s if isinstance(s, list) else [(s, {})]
    first = True
    for run_text, style in runs:
        for i, line in enumerate(str(run_text).split("\n")):
            if first:
                p = tf.paragraphs[0]
                first = False
            elif i > 0 or style.get("newline", True):
                p = tf.add_paragraph()
            p.alignment = align
            if spacing:
                p.space_after = Pt(spacing)
            r = p.add_run()
            r.text = line
            f = r.font
            f.name = style.get("font", font)
            f.size = Pt(style.get("size", size))
            f.bold = style.get("bold", bold)
            f.color.rgb = style.get("color", color)
    return tb


def new_slide(dark=True):
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide


def footer(slide, number=True):
    global page_no
    page_no += 1
    text(slide, LM, H - Inches(0.42), Inches(8), Inches(0.3), FOOTER, size=8, color=DIM)
    if number:
        text(slide, W - LM - Inches(1), H - Inches(0.42), Inches(1), Inches(0.3), str(page_no), size=9,
             color=MUTED, align=PP_ALIGN.RIGHT)


def header(slide, title, subtitle):
    text(slide, LM, Inches(0.35), CW, Inches(0.8), title, size=32, bold=True)
    text(slide, LM, Inches(1.08), CW, Inches(0.5), subtitle, size=17, color=ACCENT)


def callout(slide, y, label, body, h=Inches(0.95)):
    rect(slide, LM, y, CW, h, fill=PANEL, line=None, radius=False)
    bar = rect(slide, LM, y, Inches(0.06), h, fill=ACCENT, line=None, radius=False)
    text(slide, LM + Inches(0.25), y + Inches(0.08), CW - Inches(0.4), Inches(0.3), label, size=11, color=ACCENT, bold=True)
    text(slide, LM + Inches(0.25), y + Inches(0.36), CW - Inches(0.4), h - Inches(0.4), body, size=12, color=WHITE)


def cards(slide, y, items, h=Inches(2.05), gap=Inches(0.18), label_color=ACCENT, body_size=11):
    """items: [(label, head, body)]"""
    n = len(items)
    cw = (CW - gap * (n - 1)) / n
    for i, (label, head, body) in enumerate(items):
        x = LM + (cw + gap) * i
        rect(slide, x, y, cw, h, fill=CARD, line=BORDER)
        text(slide, x + Inches(0.18), y + Inches(0.12), cw - Inches(0.3), Inches(0.3), label, size=10, color=label_color, bold=True)
        text(slide, x + Inches(0.18), y + Inches(0.42), cw - Inches(0.3), Inches(0.45), head, size=15, bold=True)
        ln = slide.shapes.add_connector(1, x + Inches(0.18), y + Inches(0.95), x + cw - Inches(0.18), y + Inches(0.95))
        ln.line.color.rgb = BORDER
        ln.line.width = Pt(0.75)
        text(slide, x + Inches(0.18), y + Inches(1.05), cw - Inches(0.3), h - Inches(1.1), body, size=body_size, color=MUTED)


def stats(slide, items, y=None):
    """하단 통계 4~5개: [(big, small)]"""
    y = y if y is not None else H - Inches(1.45)
    ln = slide.shapes.add_connector(1, LM, y - Inches(0.1), LM + CW, y - Inches(0.1))
    ln.line.color.rgb = BORDER
    ln.line.width = Pt(0.75)
    n = len(items)
    cw = CW / n
    for i, (big, small) in enumerate(items):
        x = LM + cw * i
        text(slide, x, y, cw - Inches(0.1), Inches(0.4), big, size=17, color=ACCENT, bold=True)
        text(slide, x, y + Inches(0.4), cw - Inches(0.1), Inches(0.35), small, size=11, color=MUTED)


def matrix(slide, y, columns, rows, col_label_w=Inches(1.6), row_h=Inches(0.62), head_h=Inches(0.62),
           size=10.5, gap=Inches(0.1)):
    """columns: [(name, sub)], rows: [(label, [cell,...])]"""
    n = len(columns)
    cw = (CW - col_label_w - gap * n) / n
    for j, (name, sub) in enumerate(columns):
        x = LM + col_label_w + gap + (cw + gap) * j
        rect(slide, x, y, cw, head_h, fill=CARD, line=ACCENT, lw=1.25)
        text(slide, x, y + Inches(0.05), cw, Inches(0.32), name, size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        if sub:
            text(slide, x, y + Inches(0.33), cw, Inches(0.28), sub, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    for i, (label, cells) in enumerate(rows):
        ry = y + head_h + gap + (row_h + gap) * i
        text(slide, LM, ry, col_label_w, row_h, label, size=11, color=ACCENT, bold=True, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.RIGHT)
        for j, cell in enumerate(cells):
            x = LM + col_label_w + gap + (cw + gap) * j
            rect(slide, x, ry, cw, row_h, fill=CARD, line=BORDER)
            text(slide, x + Inches(0.08), ry, cw - Inches(0.16), row_h, cell, size=size, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


def flow(slide, y, steps, h=Inches(1.55), highlight_idx=(), gap=Inches(0.28)):
    """steps: [(num_or_None, name, layer, desc)] — 가로 화살표 흐름."""
    n = len(steps)
    cw = (CW - gap * (n - 1)) / n
    for i, (num, name, layer, desc) in enumerate(steps):
        x = LM + (cw + gap) * i
        hl = i in highlight_idx
        rect(slide, x, y, cw, h, fill=CARD, line=ACCENT if hl else BORDER, lw=1.5 if hl else 0.75)
        yy = y + Inches(0.1)
        if num:
            text(slide, x, yy, cw, Inches(0.35), str(num), size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
            yy += Inches(0.35)
        text(slide, x, yy, cw, Inches(0.35), name, size=13, bold=True, align=PP_ALIGN.CENTER)
        if layer:
            text(slide, x, yy + Inches(0.33), cw, Inches(0.25), layer, size=9, color=ACCENT, align=PP_ALIGN.CENTER)
        text(slide, x + Inches(0.06), yy + Inches(0.6), cw - Inches(0.12), h - Inches(0.8), desc, size=9.5, color=MUTED,
             align=PP_ALIGN.CENTER)
        if i < n - 1:
            ax = x + cw + Inches(0.03)
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, y + h / 2 - Inches(0.08), gap - Inches(0.06), Inches(0.16))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MUTED
            arrow.line.fill.background()


def code_block(slide, x, y, w, h, title, code):
    rect(slide, x, y, w, h, fill=CARD, line=BORDER)
    text(slide, x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Inches(0.3), title, size=10, color=ACCENT, bold=True, font=MONO)
    text(slide, x + Inches(0.15), y + Inches(0.42), w - Inches(0.3), h - Inches(0.5), code, size=9.5, color=WHITE, font=MONO)


def kv_rows(slide, x, y, w, items, row_h=Inches(0.5), gap=Inches(0.1), key_w=Inches(1.7), key_mono=True):
    for i, (k, v) in enumerate(items):
        ry = y + (row_h + gap) * i
        rect(slide, x, ry, w, row_h, fill=CARD, line=BORDER)
        text(slide, x + Inches(0.15), ry, key_w, row_h, k, size=10.5, color=ACCENT, bold=True, anchor=MSO_ANCHOR.MIDDLE,
             font=MONO if key_mono else FONT)
        text(slide, x + key_w + Inches(0.2), ry, w - key_w - Inches(0.3), row_h, v, size=10.5, anchor=MSO_ANCHOR.MIDDLE)


def bullets(slide, x, y, w, h, title, items, title_color=ACCENT, size=11):
    rect(slide, x, y, w, h, fill=CARD, line=BORDER)
    text(slide, x + Inches(0.18), y + Inches(0.12), w - Inches(0.3), Inches(0.35), title, size=13, color=title_color, bold=True)
    body = "\n".join(f"•  {it}" for it in items)
    text(slide, x + Inches(0.18), y + Inches(0.55), w - Inches(0.3), h - Inches(0.6), body, size=size, color=WHITE, spacing=4)


# ── 슬라이드 타입 ────────────────────────────────────────────────────────────
def slide_title(kicker, title, subtitle, lines, meta):
    s = new_slide()
    # 우측 장식: 액센트 세로 밴드 + 격자
    rect(s, W - Inches(1.6), 0, Inches(1.6), H, fill=PANEL, line=None, radius=False)
    for k in range(9):
        rect(s, W - Inches(1.25), Inches(0.5) + Inches(0.72) * k, Inches(0.9), Inches(0.42),
             fill=ACCENT if k in (3, 4) else CARD, line=None)
    rect(s, LM, Inches(0.45), Inches(3.4), Inches(0.36), fill=ACCENT_DIM, line=None, radius=False)
    text(s, LM + Inches(0.12), Inches(0.47), Inches(3.3), Inches(0.32), kicker, size=10, color=WHITE, bold=True)
    text(s, LM, Inches(2.2), Inches(10.5), Inches(1.1), title, size=48, bold=True)
    text(s, LM, Inches(3.3), Inches(10.5), Inches(0.6), subtitle, size=20, color=ACCENT)
    text(s, LM, Inches(4.05), Inches(10.5), Inches(0.9), lines, size=13, color=WHITE, spacing=4)
    text(s, LM, Inches(5.5), Inches(10), Inches(1.2), meta, size=13, color=MUTED, spacing=3)
    footer(s, number=False)


def slide_agenda(rows):
    s = new_slide()
    header(s, "Agenda", "11개 섹션으로 분해한 evopi 아키텍처 — Claude Code 8레이어 맵과 같은 순서로")
    y0 = Inches(1.7)
    rh = Inches(0.44)
    for i, (num, title, desc) in enumerate(rows):
        y = y0 + rh * i
        rect(s, LM, y, CW, rh - Inches(0.05), fill=PANEL, line=None, radius=False)
        text(s, LM + Inches(0.2), y, Inches(0.7), rh, num, size=13, color=ACCENT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        text(s, LM + Inches(1.0), y, Inches(4.5), rh, title, size=13, anchor=MSO_ANCHOR.MIDDLE)
        text(s, LM + Inches(5.4), y, Inches(6.5), rh, desc, size=11, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    footer(s)


def slide_divider(num, title, sub=""):
    s = new_slide()
    # 우측 기하 장식 (그라데이션 대체)
    rect(s, Inches(7.8), Inches(-0.5), Inches(6.5), Inches(8.5), fill=PANEL, line=None, radius=False)
    for k in range(6):
        rect(s, Inches(8.4) + Inches(0.75) * k, Inches(1.2) + Inches(0.55) * k, Inches(3.6), Inches(0.3),
             fill=ACCENT if k == 2 else CARD, line=None)
    text(s, LM, Inches(2.3), Inches(6), Inches(0.8), num, size=30, color=ACCENT, bold=True)
    text(s, LM, Inches(3.15), Inches(7.2), Inches(1.4), title, size=40, bold=True)
    if sub:
        text(s, LM, Inches(4.5), Inches(7.2), Inches(0.8), sub, size=15, color=MUTED)
    footer(s, number=False)


def slide_content(title, subtitle, co=None, cds=None, st=None, card_h=None, card_y=None, body_size=11):
    s = new_slide()
    header(s, title, subtitle)
    y = Inches(1.7)
    if co:
        callout(s, y, *co)
        y += Inches(1.15)
    if cds:
        ch = card_h or (Inches(2.55) if st else Inches(3.9))
        cards(s, card_y or y, cds, h=ch, body_size=body_size)
    if st:
        stats(s, st)
    footer(s)
    return s


def slide_image(title, subtitle, path, caption="", top=Inches(1.6)):
    s = new_slide()
    header(s, title, subtitle)
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    max_w, max_h = CW, H - top - Inches(0.9)
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    x = int(LM + (CW - w) / 2)
    s.shapes.add_picture(path, x, top, width=w, height=h)
    if caption:
        text(s, LM, H - Inches(0.85), CW, Inches(0.35), caption, size=10, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s)


def slide_matrix(title, subtitle, columns, rows, st=None, **kw):
    s = new_slide()
    header(s, title, subtitle)
    matrix(s, Inches(1.7), columns, rows, **kw)
    if st:
        stats(s, st)
    footer(s)
    return s


def slide_summary(title, subtitle, items, st):
    s = new_slide()
    header(s, title, subtitle)
    n = len(items)
    cols = 5 if n > 8 else 4
    gap = Inches(0.15)
    cw = (CW - gap * (cols - 1)) / cols
    ch = Inches(1.6)
    for i, (label, head, body) in enumerate(items):
        r, c = divmod(i, cols)
        x = LM + (cw + gap) * c
        y = Inches(1.7) + (ch + gap) * r
        hl = label == "EVO"
        rect(s, x, y, cw, ch, fill=CARD, line=ACCENT if hl else BORDER, lw=1.5 if hl else 0.75)
        text(s, x + Inches(0.15), y + Inches(0.1), cw - Inches(0.3), Inches(0.3), label, size=10, color=ACCENT, bold=True)
        text(s, x + Inches(0.15), y + Inches(0.4), cw - Inches(0.3), Inches(0.4), head, size=14, bold=True)
        text(s, x + Inches(0.15), y + Inches(0.85), cw - Inches(0.3), ch - Inches(0.9), body, size=9.5, color=MUTED)
    stats(s, st)
    footer(s)


def slide_refs(items):
    s = new_slide()
    header(s, "References", "공식 문서·원본 레포·논문·본 프로젝트 산출물")
    y = Inches(1.7)
    for label, head, sub in items:
        rect(s, LM, y, CW, Inches(0.78), fill=PANEL, line=None, radius=False)
        rect(s, LM, y, Inches(0.06), Inches(0.78), fill=ACCENT, line=None, radius=False)
        text(s, LM + Inches(0.25), y + Inches(0.06), CW, Inches(0.28), label, size=10, color=ACCENT, bold=True)
        text(s, LM + Inches(0.25), y + Inches(0.3), CW, Inches(0.28), head, size=13, bold=True)
        text(s, LM + Inches(0.25), y + Inches(0.53), CW, Inches(0.25), sub, size=9.5, color=MUTED)
        y += Inches(0.9)
    footer(s)


def slide_thanks():
    s = new_slide()
    rect(s, Inches(8.6), Inches(-0.5), Inches(5.5), Inches(8.5), fill=PANEL, line=None, radius=False)
    for k in range(7):
        rect(s, Inches(9.2), Inches(0.6) + Inches(0.9) * k, Inches(3.8), Inches(0.35),
             fill=ACCENT if k == 3 else CARD, line=None)
    text(s, LM, Inches(2.5), Inches(7.5), Inches(1.2), "Thank you!", size=48, bold=True)
    text(s, LM, Inches(4.2), Inches(7.5), Inches(1.6),
         "oh-my-evopi · https://github.com/sunwoo95/oh-my-evopi\n"
         "설치: curl -fsSL https://sunwoo95.github.io/oh-my-evopi/install.sh | sh\n"
         "실행: evopi   ·   evo 레이어 토글: EVOPI_EVO=on|off", size=13, color=MUTED, spacing=4)
    footer(s, number=False)


# ═════════════════════════════════════════════════════════════════════════════
# 슬라이드 시나리오
# ═════════════════════════════════════════════════════════════════════════════

# 1 표지
slide_title(
    "EVOPI ARCHITECTURE 이해",
    "evopi Architecture",
    "자기진화 하네스 구조 — Claude Code · oh-my-pi · prime-agent 와 비교",
    "레포 — https://github.com/sunwoo95/oh-my-evopi\n설치 — curl -fsSL https://sunwoo95.github.io/oh-my-evopi/install.sh | sh",
    "oh-my-evopi 팀 세미나 · 2026-09\n근거: docs/analysis/*.md, docs/design/DECISIONS.md (파일:라인 인용)\n"
    "참조 시나리오: Claude Code Architecture (AWS Korea 최우형, 2026)",
)

# 2 Agenda
slide_agenda([
    ("01", "evopi 개요", "prime 골격 + omp 자산 + Evo-Harness 델타 — 무엇을 왜 만들었나"),
    ("02", "전체 아키텍처 맵", "Master Agent Loop 중심 8 레이어 + evo 레이어, Claude Code 맵과 대응"),
    ("03", "The Agentic Loop", "Think → ipython 실행 → Verify, agentLoop + AgentSession"),
    ("04", "Models & Tools", "9 API 종 · 단일 툴 ipython · IPython 커널(uv + dill)"),
    ("05", "에이전트 구동 + Skills", "5가지 입력 파일, Python 스킬 계약, Harness 원장"),
    ("06", "Input Layer", "인터페이스 6종, curl 설치, 세션(jsonl+dill), 권한 게이트"),
    ("07", "Knowledge Layer", "시스템 프롬프트 조립, 하네스 주입(MMR), 요약 컴팩션"),
    ("08", "Execution & Multi-Agent", "sdk streamFn 파이프라인(auth-pool·dialect), rlm 서브에이전트"),
    ("09", "Observability & Integration", "확장 훅 31종, 빌트인 3종, MCP(커널 내), 프로바이더, natives"),
    ("10", "Evo Layer", "논문 ↔ prime 델타 판정, grounded-refine(D1+D4), 4-arm 평가"),
    ("11", "비교와 점검", "4자 비교 매트릭스, 하네스 점검(강점·리스크), End-to-End, 요약"),
])

# ── 01 개요 ──────────────────────────────────────────────────────────────────
slide_divider("01", "evopi 개요", "prime-agent 골격 · oh-my-pi 자산 · Evo-Harness 델타")
slide_content(
    "evopi 개요", "스스로 하네스를 다듬는 코딩 에이전트 CLI",
    co=("SELF-EVOLVING HARNESS",
        "evopi 는 prime-agent 의 골격(RLM 하네스 + IPython 커널)을 그대로 쓰고, oh-my-pi 의 TypeScript 자산"
        "(dialect · auth-pool · hashline · mnemopi)을 선별 이식한 뒤, Evo-Harness 논문(arXiv 2608.15071)의 "
        "접지 피드백 델타를 optional 레이어로 얹은 코딩 에이전트다. 설정 경로는 ~/.evopi 하나, 실행 커맨드는 evopi."),
    cds=[
        ("RLM-NATIVE", "파이썬 커널이 곧 도구",
         "내장 LLM 툴은 ipython 하나. 셸·편집·서브에이전트·MCP 를 커널 안 Python 호출(rlm.bash, edit 스킬, rlm())로 합성한다."),
        ("SELF-EVOLVING", "하네스 원장을 스스로 편집",
         "prompt / memory / skill / subagent 4종 엔트리를 refine 라운드가 create·update·delete. 롤백 스냅샷과 refinements.jsonl 이력."),
        ("PROVIDER-AGNOSTIC", "어느 모델이든 같은 루프",
         "9 API 종 · OAuth 3종 · Bedrock · Databricks · 오픈모델 in-band 툴콜(dialect 11종) · 다중 키 풀 로테이션."),
        ("EVO OPTIONAL", "evo off = prime 원본",
         "EVOPI_EVO=off 면 grounded-refine·MMR 선택기가 로드되지 않아 prime 동작과 바이트 동일 — A/B 대조군이 공짜로 생긴다."),
    ],
    st=[("1 tool", "내장 LLM 툴 = ipython"), ("4 kinds", "하네스 엔트리 종류"), ("31 hooks", "확장 훅 이벤트"),
        ("4 arms", "A/B 평가 조건")],
)

# ── 02 전체 아키텍처 맵 ──────────────────────────────────────────────────────
slide_divider("02", "전체 아키텍처 맵", "Master Agent Loop 중심 8 레이어 + evo 레이어")
slide_image("전체 아키텍처 맵", "Master Agent Loop(prime 골격)을 중심으로 한 여덟 개 레이어 + evo 레이어",
            os.path.join(DIAG, "evopi-master-arch.png"),
            "docs/diagrams/evopi-master-arch.dot (graphviz) — 각 박스 부제는 코드 실측(docs/analysis/evopi-harness-inventory.md)")
slide_image("참조: Claude Code 전체 아키텍처 맵", "같은 8 레이어 격자 — evopi 맵과 박스 단위로 대응된다",
            os.path.join(DIAG, "claude-code-master-arch.png"),
            "참조 슬라이드(AWS Korea, Claude Code Architecture p.3) 전사본을 동일 스타일로 재현 — refs/claude-code-8layer-map.md")
slide_matrix(
    "맵 대응표 — Claude Code ↔ evopi", "박스 하나하나가 어디로 대응되는가, 그리고 evopi 만의 9번째 레이어",
    [("Claude Code", "참조 맵"), ("evopi", "대응 구성요소"), ("차이의 본질", "")],
    [
        ("INPUT", ["CLI/IDE/CI-CD · Resume/Fork · Ask/Allow/Deny",
                   "TUI/print/RPC/ACP/SDK/daemon · jsonl 트리+dill 복원 · permission-gate block/warn/off",
                   "세션 복원이 대화 + 커널 변수(dill)까지 포함"]),
        ("KNOWLEDGE", ["CLAUDE.md · Auto Memory · Skills · Compaction",
                       "AGENTS.md/SYSTEM.md · Harness Ledger · SKILL.md+pyproject · 요약 컴팩션+MMR",
                       "메모리가 파일이 아니라 편집 가능한 구조화 원장"]),
        ("EXECUTION", ["Tool Dispatch(typed) · Prompt Cache · Streaming",
                       "ipython 단일 툴(sequential) · IPython 커널 · streamFn(auth-pool·dialect)",
                       "다수 도구 → 단일 REPL 셀로 합성"]),
        ("MULTI-AGENT", ["Subagents(isolated) · Worktrees",
                         "rlm() 커널 내 호출 → in-process 자식 Agent · Worktrees 는 v2",
                         "서브에이전트가 Python 함수 호출"]),
        ("OBSERVABILITY", ["Hooks(lifecycle) · Background", "확장 훅 31 이벤트 · rlm.bash 핸들/cron/백그라운드 refine",
                           "훅이 refine 플래너까지 교체 가능(session_before_refine)"]),
        ("INTEGRATION", ["MCP Runtime · Ext Servers", "rlm.mcp(커널 내) · 9 API 프로바이더 · pi-natives",
                         "MCP 도구도 Python 네임스페이스로 노출"]),
        ("OUTPUT", ["Task Result / Memory updated", "Task Result / harness 갱신(refinements.jsonl)",
                    "출력이 곧 다음 하네스 편집의 증거"]),
        ("(+) EVO", ["— (없음)", "autoRefine + grounded-refine(D1·D4) · EVOPI_EVO 게이트 · metaharness 접지",
                     "Claude Code 에 없는 9번째 레이어"]),
    ],
    row_h=Inches(0.5), head_h=Inches(0.5), size=9.5, col_label_w=Inches(1.4),
)

# ── 03 Agent Loop ────────────────────────────────────────────────────────────
slide_divider("03", "The Agentic Loop", "Think → ipython 실행 → Verify")
s = new_slide()
header(s, "The Agentic Loop", "모든 작업을 관통하는 루프 — 도구가 하나라서 루프가 단순하다")
callout(s, Inches(1.7), "CORE LOOP",
        "agentLoop()(packages/agent/src/agent-loop.ts, 963줄)이 turn_start → 모델 스트리밍 → 툴 실행 → turn_end 를 반복한다. "
        "AgentSession(agent-session.ts ≈12k줄)이 그 바깥에서 시스템 프롬프트 재조립, 자동 컴팩션, refine 체크포인트를 감싼다.")
flow(s, Inches(3.05), [
    (None, "Prompt", "Input", "사용자 요청 · steering / follow-up 큐"),
    (None, "Think", "모델 스트림", "streamFn → tool_use(ipython) 생성"),
    (None, "ipython 실행", "Execution", "셸·편집·테스트·서브에이전트를 한 셀에"),
    (None, "Verify", "커널 출력", "stdout/stderr/traceback 이 다음 턴 입력"),
    (None, "Complete", "Output", "결과 + (evo) 하네스 갱신"),
], h=Inches(1.3), highlight_idx=(0, 4))
cards(s, Inches(4.5), [
    ("질문형 작업", "커널 없이 텍스트만", "ipython 호출이 없으면 순수 스트리밍 응답으로 끝난다."),
    ("코드 수정", "edit 스킬 → 테스트 반복", "Python edit 스킬로 수정 후 같은 셀에서 pytest/tsc 실행, 실패면 재시도."),
    ("장기 작업", "daemon goal / cron / heartbeat", "데몬 슈퍼바이저가 세션 워커를 유지, 백그라운드 refine 플랜 병행."),
], h=Inches(1.35), body_size=10)
stats(s, [("steering", "턴 중 지시 주입"), ("follow-up", "턴 종료 후 후속"), ("compact", "자동 압축 → refine 트리거"),
          ("sequential", "커널 단일 스레드 — 툴 병렬 없음")])
footer(s)

# ── 04 Models & Tools ────────────────────────────────────────────────────────
slide_divider("04", "Models & Tools", "추론하는 모델 · 행동하는 단일 툴 · 그 뒤의 커널")
s = new_slide()
header(s, "Models & Tools", "루프를 움직이는 두 축 — 9 API 종의 모델과 ipython 하나의 도구")
half = (CW - Inches(0.2)) / 2
rect(s, LM, Inches(1.7), half, Inches(0.95), fill=PANEL, line=None, radius=False)
rect(s, LM, Inches(1.7), Inches(0.06), Inches(0.95), fill=ACCENT, line=None, radius=False)
text(s, LM + Inches(0.25), Inches(1.78), half, Inches(0.3), "MODELS — 추론 엔진", size=11, color=ACCENT, bold=True)
text(s, LM + Inches(0.25), Inches(2.06), half - Inches(0.4), Inches(0.6),
     "pi-ai stream(): anthropic-messages · openai-completions/responses · azure · codex · google · vertex · mistral · bedrock. "
     "models.generated 카탈로그 + models.json 사용자 등록. OAuth 3종, Databricks, prime-inference 는 peer.", size=10.5)
x2 = LM + half + Inches(0.2)
rect(s, x2, Inches(1.7), half, Inches(0.95), fill=PANEL, line=None, radius=False)
rect(s, x2, Inches(1.7), Inches(0.06), Inches(0.95), fill=ACCENT, line=None, radius=False)
text(s, x2 + Inches(0.25), Inches(1.78), half, Inches(0.3), "TOOLS — 행동 수단", size=11, color=ACCENT, bold=True)
text(s, x2 + Inches(0.25), Inches(2.06), half - Inches(0.4), Inches(0.6),
     "레지스트리 등록 툴은 ipython(기본)과 hashline_edit(선택) 둘뿐(tools/index.ts:59). TS bash/edit 정의는 있지만 미등록 — "
     "셸과 편집은 커널 안 Python(rlm.bash, edit 스킬)이 담당한다.", size=10.5)
cards(s, Inches(2.85), [
    ("IPYTHON", "영속 REPL", "단일 문자열 code 파라미터. sequential 실행. 65,536자 출력 캡."),
    ("BASH", "rlm.bash()", "프로세스 그룹 저널링, 핸들(tail/poll/kill), 커널 종료 시 회수."),
    ("EDIT", "edit 스킬 · hashline_edit", "Python edit 스킬이 기본. hashline 앵커 편집은 --tools 게이트."),
    ("SUBAGENT", "rlm()", "커널에서 함수 호출 → host_request → in-process 자식 Agent."),
    ("MCP", "rlm.mcp", "MCP 서버 도구를 Python 네임스페이스로 노출. settings 로 등록."),
], h=Inches(2.05), body_size=10)
stats(s, [("9", "API 종 (register-builtins)"), ("1 + 1", "내장 툴 ipython + 선택 hashline_edit"),
          ("11", "in-band 툴콜 dialect"), ("N keys", "EVOPI_API_KEY_POOL_<PROVIDER>")])
footer(s)

slide_content(
    "IPython Kernel", "prime 에서 무변경 이식한 실행 엔진 — 세션을 넘어 살아있는 네임스페이스",
    co=("KERNEL CONTRACT",
        "python -m rlm.repl 서브프로세스. stdin 으로 JSON-lines 요청, stdout 으로 이벤트(ready/stdout/stderr/result/display/host_request/error/done), "
        "프로토콜 v3 정확 일치. 모든 요청은 단일 큐로 직렬화된다(repl-manager.ts:769). 커널 = 단일 스레드."),
    cds=[
        ("BOOTSTRAP", "uv 부트스트랩",
         "uv python install 3.11 → venv → runtime + dill + 12 패키지(전부 wheel, R5) → Python 스킬 editable 설치. "
         "런타임 소스 해시가 바뀌면 venv 자동 재빌드(bootstrap.ts)."),
        ("REPL MANAGER", "수명주기 1,502줄",
         "spawn → ready 30s → 셀 실행 → interrupt(1s 후 호스트 측 강제 abort) → shutdown(5s). 손상 프레임은 SIGKILL 후 스냅샷 복원으로 복구."),
        ("SNAPSHOT", "dill per-variable",
         "성공 실행 1.5s 후 kernel-state.dill/.json 저장(256MB/16MB 상한, 언피클 불가 이름은 skip 보고). resume 시 복원 → 모델은 변수를 계속 본다."),
        ("BOOT GATE", "부팅 세마포어",
         "동시 커널 부팅 min(16, max(4, cpu×2)) 로 제한(EVOPI_MAX_CONCURRENT_KERNEL_BOOTS). 실행 리소스 제한은 아니다."),
    ],
    st=[("30s", "ready 타임아웃"), ("65,536", "셀 출력 문자 캡"), ("256MB", "스냅샷 총량 상한"),
        ("비격리", "uid/rlimit/ns 없음 — D3 폴백")],
    body_size=10,
)

# ── 05 에이전트 구동 + Skills ────────────────────────────────────────────────
slide_divider("05", "에이전트 구동 및 Skills", "5가지 입력 · Python 스킬 계약 · Harness 원장")
s = new_slide()
header(s, "에이전트를 움직이는 5가지 입력", "AGENTS.md / SKILL.md / Prompt / Hook — 그리고 Claude Code 에 없는 다섯 번째, Harness 원장")
callout(s, Inches(1.7), "FIVE INPUTS",
        "Claude Code 의 4파일(CLAUDE.md·SKILL.md·Prompt·Hook)이 그대로 대응되고, 여기에 에이전트 스스로 편집하는 "
        "harness_state.json 이 추가된다. 사람이 쓰는 지식(AGENTS.md)과 에이전트가 쓰는 지식(Harness)이 분리된다.")
cards(s, Inches(2.95), [
    ("MEMORY", "AGENTS.md / SYSTEM.md", "~/.evopi/agent 와 프로젝트 .evopi/agent. 매 턴 시스템 프롬프트에 포함. APPEND_SYSTEM.md 로 덧붙임."),
    ("SKILL", "SKILL.md (+pyproject)", "md 스킬은 설명 매칭 시 본문 로드. Python 스킬은 커널에 import 되어 함수로 호출된다."),
    ("PROMPT", "Prompt", "지금의 일회성 지시. steering 으로 턴 중간에도 주입 가능."),
    ("HOOK", "Extension", "jiti 로 무컴파일 TS 로드. 31 이벤트에 핸들러. tool_call 차단, 플래너 교체까지."),
    ("HARNESS", "harness_state.json", "prompt/memory/skill/subagent 엔트리. refine 라운드가 편집, 시스템 프롬프트에 kind 별 요약 주입."),
], h=Inches(2.55), body_size=9.5)
stats(s, [("항상", "AGENTS.md — 매 턴"), ("매칭·import", "SKILL.md — 설명 / 커널"), ("지금", "Prompt"),
          ("이벤트", "Hook — 31 종"), ("자동 진화", "Harness — refine 라운드")])
footer(s)

slide_matrix(
    "무엇이 언제 컨텍스트에 들어오나", "로딩 시점 · 트리거 · 갱신 주체 비교",
    [("AGENTS.md", "항상"), ("SKILL.md", "필요할 때"), ("Prompt", "지금"), ("Hook", "자동"), ("Harness", "자동 진화")],
    [
        ("로드 시점", ["세션 시작 + 매 턴 재조립", "md: 설명 매칭 시 본문\npy: 커널 부팅 시 import", "사용자 입력 시점", "라이프사이클 이벤트", "매 턴 요약(kind별 6개) 주입"]),
        ("트리거", ["자동(조건 없음)", "description 매칭 / 코드 호출", "수동(사람)", "tool_call, turn_end 등", "turn_interval 25 · compact · 외부 실패 신호"]),
        ("갱신 주체", ["사람", "사람 (+skill-creator 스킬)", "사람", "사람(개발자)", "에이전트(refine) + 롤백"]),
        ("대표 예시", ["코딩 규칙, 빌드 명령", "edit · websearch · refine", "\"이 버그 고쳐줘\"", "permission-gate 차단", "\"이 레포는 pnpm 사용\" 메모리"]),
    ],
    row_h=Inches(0.72), size=10,
    st=[("~/.evopi/agent", "전역 설정 루트"), (".evopi/agent", "프로젝트 오버레이"), ("~/.agents/skills", "벤더중립 스킬 경로 유지"),
        ("harness/", "전역 + 세션 로컬 원장")],
)

s = new_slide()
header(s, "Skills — Python 스킬 계약", "SKILL.md 오픈 표준 위에 pyproject 를 얹어 커널이 import 하는 스킬")
callout(s, Inches(1.7), "SKILL CONTRACT",
        "SKILL.md(프론트매터 name/description) + pyproject.toml + src/<import_name>/__init__.py. 스킬 로더가 SKILL.md 를 찾고 "
        "pyproject 가 있으면 kind:\"python\" 으로 분류 → uv 가 editable 설치 → 부트스트랩 코드가 import. 설치 실패는 경고만(비치명).")
code_block(s, LM, Inches(2.95), Inches(5.6), Inches(2.6), "packages/coding-agent/skills/edit/",
           "SKILL.md              (필수: 설명·사용법)\npyproject.toml        (name, dependencies)\nsrc/edit/__init__.py  (커널이 import)\n\n"
           "# 커널 안에서\nfrom edit import apply_patch\napply_patch(path, diff)\n\n# refine 스킬\nawait refine.run()   # host_request → TS")
kv_rows(s, LM + Inches(5.85), Inches(2.95), CW - Inches(5.85), [
    ("name", "스킬 식별자 — 프롬프트 스킬 목록에 노출"),
    ("description", "md 스킬 트리거 · 유효 YAML 필수(따옴표)"),
    ("pyproject", "의존성 + 형제 스킬 로컬 의존(위상 정렬 설치)"),
    ("hash", "pyproject 해시 일치 시 재설치 스킵"),
    ("MIME", "display: application/vnd.evopi.{diff,attachment,agent-message}"),
], row_h=Inches(0.44))
stats(s, [("16", "번들 스킬"), ("11 / 5", "Python / Markdown"), ("editable", "uv pip install -e"),
          ("3 경로", "번들 · ~/.evopi · 프로젝트 (+~/.agents)")])
footer(s)

slide_content(
    "Harness 원장 — 에이전트가 편집하는 지식", "rlm.harness: 프롬프트 노트 · 메모리 · 스킬 서술 · 서브에이전트 명세의 영속 원장",
    co=("CONTINUAL HARNESS",
        "\"두 번째 실행 엔진이 아니다\" — 하네스는 실행 코드가 아니라 다음 턴을 위한 라우팅·컨텍스트 힌트다. 세션 로컬 "
        "session-artifacts/<id>/harness/harness_state.json 과 전역 ~/.evopi/agent/harness/ 를 병합해 주입한다."),
    cds=[
        ("PROMPT", "prompt note", "시스템 프롬프트에 덧붙일 짧은 규칙. 예: \"테스트는 pnpm test 로\"."),
        ("MEMORY", "memory", "재사용 사실. scope local(세션) / global(전역). mnemopi MMR 이 다양성·최근성으로 선택."),
        ("SKILL", "skill 서술", "Python 참조(reference)를 강제 — 존재하는 스킬만 등록 가능(validateEdit)."),
        ("SUBAGENT", "subagent 명세", "rlm() 으로 띄울 특화 에이전트 정의. list/delete 는 host_request."),
    ],
    st=[("create/update/delete", "3 편집 연산 (MERGE/SKIP 없음)"), ("25 턴 · 20분", "autoRefine turnInterval · cooldown"),
        ("before/after", "롤백 스냅샷 · refinements.jsonl"), ("루트 세션만", "_rlmDepth==0 + 영속 artifact 필수")],
    body_size=10,
)

# ── 06 Input Layer ───────────────────────────────────────────────────────────
slide_divider("06", "Input Layer", "인터페이스 · 설치 · 세션 · 권한")
slide_content(
    "Input Layer — 인터페이스와 설치", "어디서 실행하고 어떻게 들어오는가",
    co=("SAME LOOP, EVERYWHERE",
        "cli.ts → cli-main.ts → main.ts 하나의 진입에서 모드만 갈린다. 설치는 curl 원라이너가 GitHub Pages 의 릴리스 tarball(SHA256 검증)을 받아 "
        "npm -g 로 설치하고, 필요 시 uv 부트스트랩까지 수행한다. 상태는 ~/.evopi/agent 한 곳."),
    cds=[
        ("TUI", "대화형 터미널", "pi-tui 기반. 스플래시(EVO 워드마크), /login · /compact · /refine 슬래시 명령."),
        ("HEADLESS", "print · RPC · ACP · SDK", "파이프/에디터 통합/프로그램 임베딩. createAgentSession() 으로 동일 세션 생성."),
        ("DAEMON", "goal · cron · heartbeat", "슈퍼바이저 + 세션 워커(AgentConnection 경계). 장기 무인 실행 — v1 동결."),
        ("INSTALL", "curl | sh", "install.sh 1,621줄: preflight → tarball+SHA256SUMS → npm -g → (옵션) uv/커널 부트스트랩. 실검증 0.9.6."),
    ],
    st=[("1 path", "~/.evopi (.omp/.prime 잔존 0)"), ("0.9.6", "게시 버전 (GitHub Pages)"), ("node ≥22.8", "런타임 (Bun 실코드 0건)"),
        ("evopi", "bin 이름 · APP_NAME")],
    body_size=10,
)
slide_content(
    "Sessions & Permissions", "세션 연속성과 2계층 권한 모델(D4)",
    co=("INTENT vs ENFORCEMENT",
        "의도 계층(permission-gate 확장)은 항상 로드되어 위험 명령을 차단한다. 집행 계층(OS 샌드박스)은 능력 프로브가 통과할 때만 활성 — "
        "현 환경은 unprivileged userns 불가로 bwrap 실패 → 컨테이너 경계가 집행 계층을 대신한다(D3 폴백)."),
    cds=[
        ("SESSIONS", "jsonl 트리 + dill", "sessions/<id>.jsonl 대화 트리, session-artifacts/<id>/ 에 kernel-state.dill · scheduled-jobs · harness. resume 시 변수까지 복원."),
        ("PERMISSION GATE", "block / warn / off", "7개 위험 패턴(rm -rf / 등) + ipython 셸 이스케이프 추출. UI 있으면 Yes/No, 없으면 즉시 block. EVOPI_PERMISSION_GATE."),
        ("SANDBOX PROBE", "bwrap 기능 테스트", "--version 이 아니라 실제 --unshare-user 실행으로 판정. 불가 시 session_start 에 경고 1회, 비활성."),
        ("PROFILES", "strict / dev / eval", "strict=모두 승인, dev=기본(block), eval=무인(컨테이너 전제 + gate off)."),
    ],
    st=[("7", "위험 명령 패턴"), ("10/10", "게이트 통합 테스트"), ("bwrap ✗", "현 환경 userns 차단"), ("컨테이너", "= 집행 계층 (문서화)")],
    body_size=10,
)

# ── 07 Knowledge Layer ───────────────────────────────────────────────────────
slide_divider("07", "Knowledge Layer", "시스템 프롬프트 조립 · 하네스 주입 · 컴팩션")
slide_content(
    "Knowledge Layer — 컨텍스트와 메모리", "무엇이 시스템 프롬프트로 들어오는가",
    co=("SYSTEM PROMPT ASSEMBLY",
        "buildSystemPrompt(system-prompt.ts) 조립 순서: RLM 프롬프트(불변) → 서브에이전트 가이드 → 하네스 상태(kind 별 6개·180자) → MCP 안내 → "
        "Additional Guidance → Project Context(AGENTS.md) → 스킬 목록 → APPEND_SYSTEM.md. evo on 이면 하네스 절단이 사전순 대신 MMR 선택기(M17)."),
    cds=[
        ("ALWAYS ON", "AGENTS.md · SYSTEM.md", "전역 + 프로젝트. _rebuildSystemPrompt 가 스킬/하네스 변경 시 재조립(6 지점)."),
        ("SKILLS", "이름 + 설명만 상주", "본문은 매칭 시. Python 스킬은 import 가능 목록으로 노출."),
        ("HARNESS", "kind 별 요약 메뉴", "+N more 로 넘침 표시. 사전순(기본) 또는 MMR(최근성 반감기 7일 × jaccard 다양성, 문자 예산)."),
        ("MEMORY", "mnemopi 병존", "prime harness memory 와 저장소 분리. v1 은 선택기(rerank)로만 소비, MCP 서버 겸용은 v2."),
    ],
    st=[("6 / kind", "표시 상한 (저장 상한 아님)"), ("180 chars", "엔트리 요약 길이"), ("MMR", "harness.selection=mmr 또는 evo on"),
        ("7일", "최근성 반감기")],
    body_size=10,
)
s = new_slide()
header(s, "Context Window 관리", "가득 차면 어떻게 되는가 — 요약 기반 컴팩션(D5)과 refine 연동")
callout(s, Inches(1.7), "AUTO-COMPACTION",
        "shouldCompact 가 한계 근접을 감지하면 findCutPoint 로 보존 구간을 정하고 generateSummary(compact 스킬)가 요약을 만든다. "
        "컴팩션 완료는 autoRefine 의 두 트리거 중 하나(reason=compact) — 압축 직전의 궤적이 하네스 편집 후보가 된다.")
flow(s, Inches(3.05), [
    (1, "한계 근접", "shouldCompact", "토큰 사용량 임계"),
    (2, "컷 포인트", "findCutPoint", "최근 턴은 원문 유지"),
    (3, "요약 생성", "generateSummary", "요약 기반 (snapcompact 는 v2)"),
    (4, "컨텍스트 교체", "compact()", "시스템 프롬프트 재조립"),
    (5, "refine 트리거", "reason=compact", "하네스 편집 후보 검토"),
], h=Inches(1.7), highlight_idx=(4,))
stats(s, [("summary", "요약 기반 압축 (D5)"), ("branch", "브랜치 요약(세션 트리)"), ("/compact", "수동 명령"),
          ("compact→refine", "압축이 진화의 신호")], y=Inches(5.3))
footer(s)

# ── 08 Execution & Multi-Agent ───────────────────────────────────────────────
slide_divider("08", "Execution & Multi-Agent", "스트림 파이프라인 · 툴 디스패치 · rlm 서브에이전트")
s = new_slide()
header(s, "Execution Layer — 스트림 파이프라인", "sdk.ts streamFn 클로저 하나에 omp 자산 두 개가 끼어든다")
callout(s, Inches(1.7), "STREAM FN",
        "prime 의 agent-loop 는 무수정. sdk.ts:288 의 streamFn 이 ① getApiKeyAndHeaders ② EVOPI_API_KEY_POOL 이 있으면 withAuthStream "
        "(401→refresh 1회+형제 전환, 403/usage-limit→전 풀 순회, replay-unsafe 경계 전만 무음 재시도) ③ streamSimple ④ dialect 가 켜지면 in-band 툴콜을 네이티브 toolcall 로 재물질화.")
flow(s, Inches(3.05), [
    (None, "Model", "tool_use 블록", "ipython(code) 생성"),
    (None, "auth-pool", "withAuthStream", "키 로테이션 · 세션 스티키(FNV)"),
    (None, "dialect", "wrapOwnedDialectStream", "11 방언 in-band → toolcall"),
    (None, "Dispatch", "executeToolCalls", "sequential (커널 단일 스레드)"),
    (None, "Kernel", "ReplKernelManager", "execute → stdout/result/display"),
    (None, "Result", "컨텍스트 환원", "다음 턴 입력"),
], h=Inches(1.3), highlight_idx=(1, 2))
cards(s, Inches(4.5), [
    ("AUTH-POOL", "다중 키 로테이션", "omp auth-retry 상태기계(a/b/c) 이식. Bun.hash → FNV-1a 순수 TS."),
    ("DIALECT", "오픈모델 툴콜", "hermes/qwen3/glm/kimi/deepseek 등. off 면 스트림 reference 동일(바이트 동일 게이트)."),
    ("NATIVES", "pi-natives prebuilt", "natives-loader 가 leaf .node 직접 require(AVX2 선택), 실패 시 TS 폴백."),
], h=Inches(1.35), body_size=10)
stats(s, [("0", "prime agent-loop 수정"), ("11 + 6", "dialect/pool-stream 테스트"), ("replay-safe", "start 이벤트까지만 재시도"),
          ("TS fallback", "natives 미가용 플랫폼")])
footer(s)
slide_content(
    "Multi-Agent Layer — rlm 서브에이전트", "서브에이전트가 Python 함수 호출이다",
    co=("DELEGATE FROM THE KERNEL",
        "모델이 셀 안에서 rlm(\"...\") 을 호출하면 host_request 가 TS 로 올라가 _createInlineRlmSubagentRuntime 이 자식 SessionManager + Agent 를 "
        "in-process 로 띄운다. 결과는 셀의 반환값으로 돌아오고, 서브에이전트 명세는 하네스 subagent kind 로 영속된다."),
    cds=[
        ("RLM()", "커널 내 호출", "rlm.run / find_models / list_subagents / delete_subagent — 4 host handler."),
        ("ISOLATION", "sub-세션 아티팩트", "session-artifacts/<id>/sub-xxxxxxxx/. 부모 컨텍스트 오염 없음(요약만 반환)."),
        ("SPEC", "subagent 엔트리", "refine 이 만든 명세를 다음 세션에서 재사용. Claude Code 의 /agents 에 대응."),
        ("WORKTREES", "v2 이연", "git worktree 병렬 격리는 미구현. omp 의 worktree 격리 이식이 백로그."),
    ],
    st=[("in-process", "자식 Agent 스폰 방식"), ("host_request", "TS↔Python 브리지"), ("1024", "dedup id 상한"), ("v2", "Worktrees")],
    body_size=10,
)

# ── 09 Observability & Integration ───────────────────────────────────────────
slide_divider("09", "Observability & Integration", "확장 훅 · 빌트인 확장 · MCP · 프로바이더")
s = new_slide()
header(s, "Observability Layer — Extension Hooks", "31개 라이프사이클 이벤트로 관찰·차단·교체")
callout(s, Inches(1.7), "EXTENSION RUNNER",
        "ExtensionRunner(runner.ts)가 이벤트를 emit, 핸들러는 jiti 로 무컴파일 로드된 TS. tool_call 은 first-block 단락(하나가 block 하면 즉시 차단), "
        "session_before_refine 은 {skip} 또는 {proposal} 로 refine 라운드를 억제하거나 플래너를 통째로 교체할 수 있다.")
kv_rows(s, LM, Inches(2.95), Inches(7.3), [
    ("session_*", "start · before_switch · before_compact · before_refine · compact · shutdown · before_tree · tree"),
    ("turn / agent", "turn_start · turn_end · before_agent_start · agent_end"),
    ("tool", "tool_call(차단 가능) · tool_result · user_bash"),
    ("model / ctx", "before_provider_request · model_select · context · input"),
], row_h=Inches(0.55), key_w=Inches(1.5))
cards_x = LM + Inches(7.55)
cw3 = CW - Inches(7.55)
rect(s, cards_x, Inches(2.95), cw3, Inches(2.5), fill=CARD, line=BORDER)
text(s, cards_x + Inches(0.18), Inches(3.05), cw3, Inches(0.3), "BUILTIN EXTENSIONS", size=10, color=ACCENT, bold=True)
text(s, cards_x + Inches(0.18), Inches(3.4), cw3 - Inches(0.3), Inches(2.0),
     "permission-gate — 항상 로드 (noExtensions 제외)\n"
     "grounded-refine — EVOPI_EVO=on 일 때만\n"
     "herdr-agent-state — HERDR_ENV=1 일 때만\n\n"
     "+ 사용자 확장: ~/.evopi/agent/extensions/, 프로젝트 .evopi/agent/extensions/\n"
     "+ Background: rlm.bash 핸들 · scheduled-jobs.json · 백그라운드 refine 플랜", size=10.5, color=WHITE, spacing=3)
stats(s, [("31", "훅 이벤트"), ("3", "빌트인 확장"), ("first-block", "tool_call 단락"), ("proposal", "플래너 교체 계약")])
footer(s)
slide_content(
    "Integration Layer — MCP · Providers · Natives", "외부 도구와 모델을 표준 계약으로 연결",
    co=("INTEGRATION IN THE KERNEL",
        "MCP 서버는 TS 가 아니라 커널 Python(rlm.mcp / mcp_base.py)이 연결하고 도구를 함수로 노출한다. 프로바이더는 pi-ai 의 열린 "
        "registerApiProvider 맵에 9 종이 등록되며, node 전용 Bedrock 은 지연 import 로 브라우저 번들 검사를 통과한다."),
    cds=[
        ("MCP", "rlm.mcp", "settings.mcp 등록 → McpManager 가 host handler 제공 → 커널이 서버 연결·도구 노출. 종료 시 서버 정리."),
        ("PROVIDERS", "9 API 종", "anthropic · openai(3) · azure · google · vertex · mistral · bedrock. Databricks 는 anthropic-messages + Bearer."),
        ("AUTH", "auth.json + OAuth", "provider 당 1 크리덴셜(auth.json) + env 풀. OAuth: Anthropic · Copilot · Codex(PKCE)."),
        ("NATIVES", "pi-natives 18.1.2", "mmrRerank · cosinePairs · vectorTopK · diffLineRuns · nodeChainAt · blockBoundaries — 6 함수 스모크 통과."),
    ],
    st=[("Python", "MCP 클라이언트 위치"), ("9", "API 종"), ("3", "OAuth 프로바이더"), ("6", "natives 함수 (R6)")],
    body_size=10,
)

# ── 10 Evo Layer ─────────────────────────────────────────────────────────────
slide_divider("10", "Evo Layer", "Evo-Harness 논문 ↔ prime refine 델타 · grounded-refine · 평가")
slide_matrix(
    "논문 ↔ prime 델타 판정 (R4)", "EVO-HARNESS(arXiv 2608.15071) 9단계 중 prime 에 없는 것만 v1 로 — 나머지는 v2",
    [("논문 단계", ""), ("prime 판정", "evo.md"), ("evopi v1 처리", "DECISIONS R4")],
    [
        ("S1 저장 포맷", ["엔트리 = trigger+규칙+evidence+scope, md+yaml", "부분 — JSON 원장, trigger/evidence 는 이벤트에만", "v2 (D6)"]),
        ("S2 Select/Inject", ["관련성 기반 선택 ≤ b, 별도 모델", "부분 — Inject 동일, Select 없음(사전순 6개)", "MMR 선택기 + 문자 예산 (M17, D8 백로그 ③ 선반영)"]),
        ("S3 컨텍스트", ["(x, τ, y, f) — 결과·피드백 슬롯", "부분 — y/f 슬롯 없음", "v1 D7 → D4 에 흡수"]),
        ("S4 Reflect", ["실패/부정 피드백 시에만", "없음 — turn_interval/compact 트리거", "v1 D1: 신호 성공이면 {skip}"]),
        ("S5 Evolver", ["ADD / MERGE / REVISE / SKIP", "부분 — create/update/delete, MERGE·SKIP 없음", "v2 (D5)"]),
        ("S6 2레벨", ["general / topic 이중 컴파일", "없음 — 단일 세션 1패스", "v2 (D2, 파일 수 경계)"]),
        ("S7 배치·예산", ["배치 16, 스킬 상한 5", "없음 — 저장 상한 0건", "v2 (D9, D5 선행 필요)"]),
        ("S8 접지", ["환경 pass/fail 필수 (Self-Generated 는 악화)", "없음 — 자가 판단 = Self-Generated 동형", "v1 D4: <external_feedback> 주입"]),
        ("S9 평가", ["No-Evolve 대조군, 3회 평균", "없음 (prime) / metaharness (omp)", "v1 D0: 4-arm, autoRefine off = 대조군"]),
    ],
    row_h=Inches(0.43), head_h=Inches(0.45), size=9, col_label_w=Inches(1.45),
)
s = new_slide()
header(s, "grounded-refine — 접지 피드백으로 진화하기", "논문 Table 4: Self-Generated 61.67 < No-Evolve 63.67 < Minimal 67.33 (SWE-bench Lite)")
callout(s, Inches(1.7), "WHY GROUNDING",
        "prime 의 refine 은 LLM 이 스스로 성공을 판단하는 Self-Generated 설정과 동형 — 논문상 No-Evolve 보다 나쁘다. evopi 는 외부 pass/fail 신호를 "
        "session_before_refine 훅에서 읽어 (a) 신호 없음 → 무개입 (b) 성공 → skip (c) 실패 → 신호를 플래너 입력에 주입한다. 접지 없는 evo-on 은 금지(SPEC §4).")
flow(s, Inches(3.05), [
    (1, "외부 러너", "metaharness", "result 이벤트 → pass/fail"),
    (2, "신호 파일", "EVOPI_FEEDBACK_FILE", "{task, status, detail?}"),
    (3, "훅 진입", "session_before_refine", "!rollbackId && hasHandlers"),
    (4, "D1 게이트", "isFailureStatus", "성공 → {skip:true}"),
    (5, "D4 주입", "<external_feedback>", "Minimal 기본 · Standard 옵트인"),
    (6, "적용", "normalize → apply", "refinements.jsonl · 롤백"),
], h=Inches(1.4), highlight_idx=(3, 4))
cards(s, Inches(4.55), [
    ("EVOPI_EVO", "on / off / 미설정", "on = 확장 로드 + autoRefine on. off = autoRefine 비활성(순수 대조군). 미설정 = prime 기본."),
    ("FALLBACK", "무음 강등 방지", "플래너 모델·키 부재 시 undefined → 내장 플래너. transient 오류는 oneshot-retry 로 재시도(M18)."),
    ("SAFETY", "SPEC §4:56", "접지 신호 미배선 evo-on arm 구성 금지. 훅 미로드 시 short-circuit → prime 경로 바이트 동일."),
], h=Inches(1.35), body_size=10)
stats(s, [("+3.7pt", "Minimal vs No-Evolve (논문)"), ("1 file", "grounded-refine.ts 208줄"), ("Minimal", "기본 세밀도 (pass/fail)"),
          ("off = prime", "D7 대조군")])
footer(s)
slide_content(
    "평가 — metaharness 4-arm A/B", "omp metaharness 를 bun 격리(eval/)로 재사용, 코딩 트랙 kind:\"edit\"",
    co=("EVAL IS OUTSIDE THE PRODUCT",
        "제품은 node 전용, 평가 러너는 구조적으로 bun 전용(import.meta.dir, Bun.which…) → eval/ 소형 bun 워크스페이스로 격리(R7). "
        "arm 은 잡네임 접두 규약(evopi-<arm>), 피실험 CLI 는 bun overrides 로 @evopi/pi-coding-agent 리포인트(Q2)."),
    cds=[
        ("evopi-omp", "상류 omp", "게시된 @oh-my-pi/pi-coding-agent 18.1.2 그대로 — 대조 상류."),
        ("evopi-prime", "prime 스켈레톤", "prime 계열 설정 대조군."),
        ("evopi-evooff", "evopi · evo off", "EVOPI_EVO=off → autoRefine 비활성. 순수 대조군(No-Evolve)."),
        ("evopi-evoon", "evopi · evo on", "EVOPI_EVO=on + EVOPI_FEEDBACK_FILE 배선 필수. 신호 = 러너 result 이벤트."),
    ],
    st=[("SKIP", "실 A/B — API 키 부재"), ("2 smoke", "mock completeSimple · D1 로직"), ("edit", "task_success · edit_success"),
        ("3회 평균", "논문 프로토콜(I.3) — 키 확보 시")],
    body_size=10,
)

# ── 11 비교와 점검 ───────────────────────────────────────────────────────────
slide_divider("11", "비교와 점검", "Claude Code · oh-my-pi · prime-agent · evopi")
CMP_COLS = [("Claude Code", "참조 자료 기준"), ("oh-my-pi", "omp v18.1.2"), ("prime-agent", "v0.9.1"), ("evopi", "v0.9.6")]
slide_matrix(
    "4자 비교 매트릭스 (1/2)", "제어 루프 · 도구/실행 · 컨텍스트 · 권한/안전",
    CMP_COLS,
    [
        ("제어 루프", ["Gather-Act-Verify\n독립 도구 병렬 실행", "agent-loop 3,010줄 · 29 툴 병렬\nsteering/aside/follow-up 3큐 · TTSR", "동일 agent-loop\nAgentSession 12k줄이 감쌈",
                    "prime 루프 무수정\nstreamFn 에 dialect·pool 삽입"]),
        ("도구 / 실행", ["5 카테고리 tool_use\n로컬/Cloud VM/Remote", "29 빌트인 툴 · hashline 편집\nRust in-process 셸(brush) · xd:// 장치", "ipython 단일 툴\nIPython 커널 + dill 스냅샷",
                     "ipython + hashline_edit(선택)\n커널 무변경 · natives prebuilt"]),
        ("컨텍스트", ["CLAUDE.md · MEMORY.md\nSkills 3단 · 자동 압축", "AGENTS.md · 메모리 백엔드 4종\n컴팩션 5전략(snapcompact PNG 포함)", "AGENTS.md · Harness 원장\n요약 컴팩션",
                    "prime + MMR 하네스 선택기\n(mnemopi 커널) · 요약 컴팩션"]),
        ("권한 / 안전", ["4 모드 · allowlist 계층\n체크포인트 되감기 · Hooks 게이트", "approval allow/deny/prompt · 3모드\nCOW worktree 격리 · OS 샌드박스 없음", "권한 프롬프트\n샌드박스는 예제 확장(bwrap)",
                     "permission-gate 내장(block/warn/off)\n+ 샌드박스 프로브(D3 폴백)"]),
    ],
    row_h=Inches(0.95), size=9.5, col_label_w=Inches(1.3),
)
slide_matrix(
    "4자 비교 매트릭스 (2/2)", "확장성 · 자기개선 · 모델 연결 · 런타임/배포 · 평가",
    CMP_COLS,
    [
        ("확장성", ["Skills · MCP · Hooks(14 이벤트)\nSubagents · Plugins", "훅 28 · 슬래시 83 · 커스텀 툴\nMCP · Advisor · 서브에이전트 task", "extensions(jiti) 31 훅\nPython 스킬 · MCP(커널) · rlm()",
                  "prime 표면 그대로\n+ 빌트인 3 확장"]),
        ("자기개선", ["MEMORY.md 갱신\n(하네스 자동 정련 없음)", "autolearn 자동 SKILL.md 저작\nlearn/manage_skill · 메모리 4백엔드", "continual harness refine\n(자가 판단 = Self-Generated)",
                   "prime refine + 접지(D4)\n+ 실패 한정(D1) + MMR 선택"]),
        ("모델 연결", ["Claude (Sonnet/Opus)\n/model 전환", "카탈로그 66 프로바이더 · dialect 11\n풀 로테이션 · gateway/broker 사이드카", "9 API · 카탈로그 32 프로바이더\nOAuth 3 · Bedrock · auth.json 1키/프로바이더",
                    "prime 카탈로그 + 풀 로테이션\n+ dialect 11 + Databricks"]),
        ("런타임 / 배포", ["(자료 미기재)\nnpm CLI · 8+ 인터페이스", "Bun 전용(Bun.* 1,282회)\nRust napi 9 crate · 단일 바이너리 설치", "node · uv Python\ninstall.sh 45KB · R2",
                       "node ≥22 · Bun 실코드 0\ncurl | sh → GitHub Pages"]),
        ("평가", ["(자료 미기재)\nHooks 로깅 · /context", "metaharness\n(harbor/edit/snapcompact)", "없음\n(autoRefine off 로 대조군만)",
                "metaharness bun 격리\n4-arm · 실 실행 SKIP(키)"]),
    ],
    row_h=Inches(0.78), size=9.5, col_label_w=Inches(1.3),
)
s = new_slide()
header(s, "evopi 하네스 점검 결과", "강점 4 · 리스크 4 — 코드 실측(docs/analysis/evopi-harness-inventory.md) 기준")
cards(s, Inches(1.6), [
    ("STRENGTH", "골격 무수정", "prime agent-loop·커널 무변경, 확장은 streamFn/훅/빌트인 seam 에만. tsgo 0 · Bun 실코드 0 · .omp 0."),
    ("STRENGTH", "대조군 내장", "EVOPI_EVO=off 가 prime 원본과 바이트 동일 — A/B 설계가 코드에 박혀 있다(D7)."),
    ("STRENGTH", "접지된 진화", "논문상 유일하게 직접 절제 근거가 있는 델타(D4)를 1파일로 배선. 안전 구속을 SPEC 에 명문화."),
    ("STRENGTH", "이식 자산 활성", "dialect·auth-pool·mnemopi·hashline·oneshot-retry 전부 소비 배선 완료(휴면 0)."),
], h=Inches(2.15), body_size=10)
cards(s, Inches(3.95), [
    ("RISK · HIGH", "커널 env 전체 상속", "repl-manager.ts:257 `...process.env` — API 키가 커널·사용자 코드에 노출(Q6 미해소). allowlist 필터 필요."),
    ("RISK · HIGH", "OS 샌드박스 미구현", "bwrap 은 프로브만, 래핑 코드는 예제 확장에만. 집행 계층 = 컨테이너 경계 전제 — 배포 문서로만 보장."),
    ("RISK · MED", "셀 타임아웃 없음", "사용자 ipython 셀에 실행 시간 상한 없음, abort 도 호스트 측 정산만 — Python 무한루프는 죽지 않는다."),
    ("RISK · MED", "실 A/B 미실행", "evo 효과 주장은 아직 논문 수치 인용. 키 확보 후 4-arm 3회 실행 전까지는 가설(GAP-4)."),
], h=Inches(2.15), body_size=10, label_color=WARN)
text(s, LM, H - Inches(1.05), CW, Inches(0.5),
     "기타: TS bash/edit 툴 정의 미등록(의도된 단일 툴 철학) · `.prime/config.json` 읽기 2곳은 외부 Prime CLI interop(승인된 예외) · "
     "RESULTS.md 의 \"pi-ai mock 없음\" 서술은 현재 providers/faux.ts 존재와 불일치(문서 갱신 필요).", size=9.5, color=MUTED)
footer(s)
s = new_slide()
header(s, "종합 플로우 — End to End", "요청에서 검증된 결과, 그리고 하네스 갱신까지 한 흐름으로")
flow(s, Inches(1.75), [
    (1, "Prompt", "Input", "사용자 요청 · steering"),
    (2, "Gate", "Permission", "위험 명령 block/warn"),
    (3, "Gather", "Knowledge", "AGENTS.md · Skills · Harness(MMR)"),
    (4, "Execute", "Execution", "ipython 셀 — bash/edit/rlm/MCP"),
    (5, "Verify", "Core Loop", "커널 출력·테스트로 결과 검증"),
    (6, "Result", "Output", "검증된 산출물"),
    (7, "Refine", "Evo", "pass/fail 접지 → 원장 편집"),
], h=Inches(1.75), highlight_idx=(0, 5, 6), gap=Inches(0.22))
rect(s, LM, Inches(3.75), CW, Inches(0.55), fill=PANEL, line=None, radius=False)
text(s, LM + Inches(0.25), Inches(3.75), Inches(3), Inches(0.55), "Observability (Hooks)", size=11, color=ACCENT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
text(s, LM + Inches(3.2), Inches(3.75), CW - Inches(3.4), Inches(0.55),
     "31 이벤트로 전 과정 관찰/검증/차단 — tool_call(2) · turn_end(5) · session_before_refine(7) · session_compact",
     size=11, anchor=MSO_ANCHOR.MIDDLE)
arrow = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, LM + Inches(0.6), Inches(4.5), CW - Inches(1.2), Inches(0.3))
arrow.fill.solid(); arrow.fill.fore_color.rgb = ACCENT_DIM; arrow.line.fill.background()
text(s, LM, Inches(4.85), CW, Inches(0.35), "미완이면 다음 셀 결정하며 루프 반복 · refine 이 편집한 하네스는 다음 세션의 Gather 로 들어온다",
     size=10.5, color=ACCENT, align=PP_ALIGN.CENTER)
stats(s, [("Input→Output", "단일 에이전틱 흐름"), ("feeds", "지식 + 하네스가 매 단계 공급"), ("hooks", "전 과정 관측/통제"),
          ("grounded", "검증 결과가 진화의 유일한 입력")])
footer(s)
slide_summary(
    "핵심 요약", "여덟 개 레이어 + 하나의 evo 레이어, 하나의 루프로 기억하기",
    [
        ("CORE", "Agentic Loop", "Think → ipython → Verify. prime agent-loop 무수정."),
        ("INPUT", "진입/세션/권한", "6 인터페이스 · jsonl+dill 복원 · permission-gate."),
        ("KNOWLEDGE", "컨텍스트/하네스", "AGENTS.md · Skills · Harness 원장(MMR) · 요약 컴팩션."),
        ("EXECUTION", "커널/스트림", "IPython 커널 · streamFn(auth-pool·dialect) · sequential."),
        ("EVO", "접지된 진화", "grounded-refine(D1+D4) · EVOPI_EVO · off = prime."),
        ("MULTI-AGENT", "rlm 서브에이전트", "커널 함수 호출 → in-process 자식 Agent. Worktrees v2."),
        ("OBSERVABILITY", "Extension Hooks", "31 이벤트 · 빌트인 3 · 플래너 교체 계약."),
        ("INTEGRATION", "MCP/Providers", "rlm.mcp(커널) · 9 API · OAuth 3 · natives 6."),
        ("OUTPUT", "Task Result", "검증된 산출물 + refinements.jsonl 갱신."),
        ("EVAL", "metaharness", "bun 격리 4-arm. 실 실행은 키 확보 시."),
    ],
    [("1 loop", "Think-Execute-Verify"), ("1 tool", "ipython — 나머지는 Python"), ("4 kinds", "하네스 원장 엔트리"),
     ("evo optional", "off = prime 원본 (D7)"), ("접지 필수", "Self-Generated 는 금지")],
)
slide_refs([
    ("PAPER", "EVO-HARNESS: Context-to-Harness Skill Compilation for Self-Evolving Agents — arXiv 2608.15071",
     "S1-S9 메커니즘, Table 4 접지 피드백 절제, Appendix E/F 프롬프트·상수 — docs/analysis/evo.md 에 라인 인용"),
    ("UPSTREAM", "prime-agent (PrimeIntellect-ai) v0.9.1 · oh-my-pi v18.1.2 · pi (earendil-works)",
     "골격 · TS 자산 공급원 · 원류. 읽기 전용 분석: docs/analysis/prime.md, omp.md, *-master-arch.md"),
    ("CLAUDE CODE", "Claude Code Architecture (AWS Korea 최우형, 2026) · code.claude.com/docs · agentskills.io · modelcontextprotocol.io",
     "8레이어 맵 참조 시나리오 — docs/analysis/claude-code-arch.md (PDF 페이지 인용, 자료 미기재 항목 분리)"),
    ("EVOPI", "github.com/sunwoo95/oh-my-evopi · docs/design/{DECISIONS,PORTING,AUDIT-initial-goal}.md · docs/analysis/evopi-harness-inventory.md",
     "결정 원장(D1-D8, R3-R10) · 이식 등급표 · 코드 실측 인벤토리 · docs/diagrams/*.dot"),
])
slide_thanks()

prs.save(OUT)
print(f"saved {OUT} ({len(prs.slides)} slides)")
